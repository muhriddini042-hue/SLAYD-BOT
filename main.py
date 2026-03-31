import os
import requests
from pptx import Presentation
from pptx.util import Inches
from telegram import Update, ReplyKeyboardMarkup
from telegram.ext import (
    ApplicationBuilder,
    CommandHandler,
    MessageHandler,
    filters,
    ContextTypes,
)
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet

# 🔐 TOKEN (Render ENV orqali beriladi)
TOKEN = os.getenv("BOT_TOKEN")

# 🧠 AI API
API_URL = "https://api.apifreellm.com/v1/chat/completions"

# 🌍 user language storage
user_lang = {}

# 🌍 til tanlash tugmalari
keyboard = [["UZ 🇺🇿", "RU 🇷🇺", "EN 🇺🇸"]]

# =========================
# 🧠 AI TEXT GENERATOR
# =========================
def generate_text(topic, lang):
    prompt = f"""
    {topic} mavzusida 5 ta slayd uchun qisqa, aniq va tushunarli matn yoz.
    Har bir slaydni alohida qatorda yoz.
    Til: {lang}
    """

    try:
        response = requests.post(API_URL, json={
            "model": "gpt-3.5-turbo",
            "messages": [{"role": "user", "content": prompt}]
        })

        data = response.json()
        text = data["choices"][0]["message"]["content"]

        slides = [line for line in text.split("\n") if line.strip() != ""]
        return slides[:5]

    except Exception as e:
        return [
            f"{topic} haqida umumiy tushuncha",
            f"{topic}ning ahamiyati",
            f"{topic} muammolari",
            f"{topic} yechimlari",
            "Xulosa"
        ]

# =========================
# 🖼 RANDOM IMAGE
# =========================
def get_image():
    return "https://picsum.photos/400"

# =========================
# 🎞 PPT GENERATOR
# =========================
def create_ppt(topic, slides):
    prs = Presentation()

    for text in slides:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = topic
        slide.placeholders[1].text = text

        # rasm qo‘shish
        img_url = get_image()
        img_data = requests.get(img_url).content

        with open("temp.jpg", "wb") as f:
            f.write(img_data)

        slide.shapes.add_picture("temp.jpg", Inches(5), Inches(2), width=Inches(3))

    file_name = f"{topic}.pptx"
    prs.save(file_name)
    return file_name

# =========================
# 📄 PDF GENERATOR
# =========================
def create_pdf(topic, slides):
    file_name = f"{topic}.pdf"
    doc = SimpleDocTemplate(file_name)
    styles = getSampleStyleSheet()

    content = []
    for slide in slides:
        content.append(Paragraph(slide, styles["Normal"]))

    doc.build(content)
    return file_name

# =========================
# ▶️ START
# =========================
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Tilni tanlang 👇",
        reply_markup=ReplyKeyboardMarkup(keyboard, resize_keyboard=True)
    )

# =========================
# 🌍 LANGUAGE SELECT
# =========================
async def choose_lang(update: Update, context: ContextTypes.DEFAULT_TYPE):
    text = update.message.text

    if text.startswith("UZ"):
        user_lang[update.effective_chat.id] = "uz"
    elif text.startswith("RU"):
        user_lang[update.effective_chat.id] = "ru"
    elif text.startswith("EN"):
        user_lang[update.effective_chat.id] = "en"
    else:
        return False

    await update.message.reply_text("Mavzuni yozing 🎯")
    return True

# =========================
# 💬 MAIN HANDLER
# =========================
async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    text = update.message.text

    # agar til tanlayotgan bo‘lsa
    lang_selected = await choose_lang(update, context)
    if lang_selected:
        return

    topic = text
    lang = user_lang.get(update.effective_chat.id, "uz")

    await update.message.reply_text("⏳ Slayd tayyorlanmoqda...")

    slides = generate_text(topic, lang)

    ppt_file = create_ppt(topic, slides)
    pdf_file = create_pdf(topic, slides)

    await update.message.reply_document(open(ppt_file, "rb"))
    await update.message.reply_document(open(pdf_file, "rb"))

# =========================
# 🚀 RUN
# =========================
def main():
    app = ApplicationBuilder().token(TOKEN).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))

    print("Bot ishga tushdi 🚀")
    app.run_polling()

if __name__ == "__main__":
    main()
